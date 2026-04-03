#!/usr/bin/env python3
from __future__ import annotations

import argparse
from pathlib import Path

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


def add_textbox(slide, x, y, w, h, text, size=12, bold=False):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.alignment = PP_ALIGN.LEFT
    return tb


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument("--company", required=True)
    parser.add_argument("--ticker", required=True)
    parser.add_argument("--out", required=True)
    parser.add_argument("--summary-out", required=True)
    args = parser.parse_args()

    out = Path(args.out)
    summary_out = Path(args.summary_out)
    out.parent.mkdir(parents=True, exist_ok=True)
    summary_out.parent.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    add_textbox(slide, 0.3, 0.2, 9.2, 0.4, f"{args.company} ({args.ticker})", size=20, bold=True)
    add_textbox(slide, 0.3, 0.8, 4.5, 2.5, "Company Overview\n- Vertical SaaS operator\n- Fictional smoke-test issuer\n- U.S.-focused platform", size=11)
    add_textbox(slide, 5.1, 0.8, 4.5, 2.5, "Business & Positioning\n- Recurring revenue model\n- Sticky workflow integration\n- Attractive upsell path", size=11)
    add_textbox(slide, 0.3, 3.7, 4.5, 2.8, "Key Financials\n- Revenue: $45.0mm\n- EBITDA Margin: 22.0%\n- Growth: Mid-teens", size=11)
    add_textbox(slide, 5.1, 3.7, 4.5, 2.8, "Developments\n- Focused GTM expansion\n- Product analytics launch\n- Sponsorable profile", size=11)

    prs.save(out)
    summary_out.write_text(
        "\n".join(
            [
                f"# {args.company} Strip Profile",
                "",
                "- Generated as a single-slide profile for smoke testing.",
                f"- Output file: {out}",
            ]
        )
        + "\n",
        encoding="utf-8",
    )


if __name__ == "__main__":
    main()
